#!/usr/bin/env python3
"""
generate_presentation.py

Generates a polished PowerPoint deck summarizing the Esports CV Pipeline project.
Light theme with high-contrast text for projector readability.

Usage:
    ./venv/bin/python3 generate_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Config ────────────────────────────────────────────────────────────
OUTPUT = "Esports_CV_Pipeline_Presentation.pptx"

# Light theme colors — high contrast for projectors
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)   # White background
BG_LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)   # Slate 100 — card background
BG_CARD    = RGBColor(0xE2, 0xE8, 0xF0)   # Slate 200 — darker card
ACCENT_GREEN = RGBColor(0x16, 0xA3, 0x4A) # Green 600 — strong green
ACCENT_RED   = RGBColor(0xDC, 0x26, 0x26) # Red 600
ACCENT_BLUE  = RGBColor(0x25, 0x63, 0xEB) # Blue 600
ACCENT_AMBER = RGBColor(0xD9, 0x77, 0x06) # Amber 600
TEXT_BLACK   = RGBColor(0x0F, 0x17, 0x2A)  # Near-black — primary text
TEXT_DARK    = RGBColor(0x1E, 0x29, 0x3B)  # Slate 800 — headings
TEXT_BODY    = RGBColor(0x33, 0x41, 0x55)  # Slate 700 — body text
TEXT_LABEL   = RGBColor(0x47, 0x55, 0x69)  # Slate 600 — labels

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Image paths
IMG_DIR = "/Users/t-dawg/Documents/esports_cv"
REFERENCE_MAP    = os.path.join(IMG_DIR, "reference_map.jpg")
SCOREBOARD_CROP  = os.path.join(IMG_DIR, "debug", "scoreboard_area.png")
VALIDATION_CHART = os.path.join(IMG_DIR, "final_validation_chart.png")
DEMO_VIDEO       = os.path.join(IMG_DIR, "demo_clip.mp4")
DEMO_THUMBNAIL   = os.path.join(IMG_DIR, "demo_thumbnail.png")


# ── Helpers ───────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       color=TEXT_BODY, line_spacing_pt=24):
    """Add multiple lines of text, each as a separate paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.line_spacing = Pt(line_spacing_pt)
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_GREEN, thickness=4):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_stat_card(slide, left, top, width, height, number, label,
                  accent=ACCENT_GREEN):
    card = add_shape_fill(slide, left, top, width, height, BG_LIGHT, corner_radius=0.05)

    # Accent bar at top of card
    add_shape_fill(slide, left + Inches(0.15), top + Inches(0.12),
                   Inches(0.6), Pt(4), accent)

    # Big number
    add_text_box(slide, left + Inches(0.15), top + Inches(0.28),
                 width - Inches(0.3), Inches(0.7),
                 number, font_size=36, color=TEXT_BLACK, bold=True)

    # Label
    add_text_box(slide, left + Inches(0.15), top + Inches(0.95),
                 width - Inches(0.3), Inches(0.5),
                 label, font_size=14, color=TEXT_BODY)


def slide_header(slide, section_label, title, accent_color=ACCENT_GREEN):
    """Standard slide header: section label + accent line + title."""
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
                 section_label, font_size=14, color=accent_color, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), accent_color)
    add_text_box(slide, Inches(0.8), Inches(1.05), Inches(11.5), Inches(0.7),
                 title, font_size=30, color=TEXT_BLACK, bold=True)


# ── Slide Builders ────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1: Title — clean, just title + author."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top accent bar
    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_GREEN)

    # Main title
    add_text_box(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1.8),
                 "Reading the Game:\nReal-Time Win Prediction from Broadcast Video",
                 font_size=42, color=TEXT_BLACK, bold=True,
                 alignment=PP_ALIGN.CENTER, line_spacing=54)

    # Divider
    add_accent_line(slide, Inches(5.5), Inches(4.3), Inches(2.3), ACCENT_GREEN, thickness=4)

    # Author
    add_text_box(slide, Inches(1.2), Inches(4.8), Inches(11), Inches(0.6),
                 "By Trevor Gregory",
                 font_size=24, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 2: The Problem — simplified."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "THE PROBLEM", "Can we predict who wins \u2014 while the match is still live?")

    # Left: What we have
    add_shape_fill(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.5),
                 "WHAT WE HAVE", font_size=16, color=ACCENT_RED, bold=True)

    add_multiline_text(slide, Inches(1.1), Inches(3.1), Inches(5.0), Inches(3.0), [
        "\u2022  Raw broadcast video only \u2014 no game API",
        "\u2022  Small, noisy scoreboard text at 60 fps",
        "\u2022  Cluttered minimap with player icons",
    ], font_size=18, color=TEXT_BODY)

    # Right: What we built
    add_shape_fill(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.5),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(7.1), Inches(2.4), Inches(5), Inches(0.5),
                 "WHAT WE BUILT", font_size=16, color=ACCENT_GREEN, bold=True)

    add_multiline_text(slide, Inches(7.1), Inches(3.1), Inches(5.0), Inches(3.0), [
        "\u2022  CV pipeline: pixels \u2192 structured data",
        "\u2022  XGBoost model: data \u2192 win probability",
        "\u2022  Broadcast overlay: probability \u2192 live video",
    ], font_size=18, color=TEXT_BODY)


def slide_pipeline(prs):
    """Slide 3: Pipeline — simplified 3-stage view."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "PIPELINE", "Three-Stage Architecture")

    stages = [
        ("1", "EXTRACT", "Scoreboard OCR\n+ Minimap Detection",
         "Template matching for scores/timer.\nYOLOv8 for player counts on minimap.",
         ACCENT_BLUE),
        ("2", "ENGINEER", "12 Features from\nTwo Data Sources",
         "6 scoreboard features (scores, rates, time)\n6 telemetry features (player counts, advantage)",
         ACCENT_AMBER),
        ("3", "PREDICT", "XGBoost Win\nProbability Model",
         "Trained on 2 matches, tested on held-out match.\nMirror augmentation for class balance.",
         ACCENT_GREEN),
    ]

    card_w = Inches(3.7)
    gap = Inches(0.35)
    start_x = Inches(0.8)
    card_y = Inches(2.2)

    for i, (num, title, subtitle, desc, accent) in enumerate(stages):
        x = start_x + i * (card_w + gap)

        # Card
        add_shape_fill(slide, x, card_y, card_w, Inches(4.5), BG_LIGHT, corner_radius=0.04)

        # Number badge
        badge = add_shape_fill(slide, x + Inches(0.2), card_y + Inches(0.25),
                               Inches(0.55), Inches(0.55), accent, corner_radius=0.15)
        add_text_box(slide, x + Inches(0.2), card_y + Inches(0.28),
                     Inches(0.55), Inches(0.5),
                     num, font_size=24, color=BG_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.9), card_y + Inches(0.3),
                     Inches(2.5), Inches(0.4),
                     title, font_size=22, color=TEXT_BLACK, bold=True)

        # Subtitle
        add_text_box(slide, x + Inches(0.2), card_y + Inches(1.1),
                     card_w - Inches(0.4), Inches(0.8),
                     subtitle, font_size=17, color=accent, bold=True,
                     line_spacing=24)

        # Description
        add_text_box(slide, x + Inches(0.2), card_y + Inches(2.2),
                     card_w - Inches(0.4), Inches(1.5),
                     desc, font_size=15, color=TEXT_BODY, line_spacing=22)

        # Arrow between stages
        if i < len(stages) - 1:
            arrow_x = x + card_w + Inches(0.02)
            add_text_box(slide, arrow_x, card_y + Inches(1.8),
                         Inches(0.3), Inches(0.5),
                         "\u25B6", font_size=20, color=TEXT_LABEL,
                         alignment=PP_ALIGN.CENTER)


def slide_cv_extraction(prs):
    """Slide 4: Computer Vision — scoreboard + minimap with game screenshot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "COMPUTER VISION", "Extracting Data from Every Frame")

    # Game screenshot (large)
    if os.path.exists(REFERENCE_MAP):
        slide.shapes.add_picture(
            REFERENCE_MAP, Inches(0.8), Inches(2.0),
            width=Inches(7.5), height=Inches(4.2),
        )

    # Right side: two extraction boxes
    # Scoreboard extraction
    add_shape_fill(slide, Inches(8.8), Inches(2.0), Inches(3.7), Inches(2.0),
                   BG_LIGHT, corner_radius=0.04)
    add_text_box(slide, Inches(9.1), Inches(2.15), Inches(3.2), Inches(0.4),
                 "SCOREBOARD OCR", font_size=16, color=ACCENT_BLUE, bold=True)
    add_multiline_text(slide, Inches(9.1), Inches(2.6), Inches(3.2), Inches(1.2), [
        "Template matching on fixed ROIs",
        "Extracts scores + timer per frame",
        "99.3% field completeness",
    ], font_size=14, color=TEXT_BODY, line_spacing_pt=20)

    # Minimap detection
    add_shape_fill(slide, Inches(8.8), Inches(4.2), Inches(3.7), Inches(2.0),
                   BG_LIGHT, corner_radius=0.04)
    add_text_box(slide, Inches(9.1), Inches(4.35), Inches(3.2), Inches(0.4),
                 "MINIMAP DETECTION", font_size=16, color=ACCENT_AMBER, bold=True)
    add_multiline_text(slide, Inches(9.1), Inches(4.8), Inches(3.2), Inches(1.2), [
        "YOLOv8 trained on minimap icons",
        "Counts friendly + enemy players",
        "Tracks player advantage over time",
    ], font_size=14, color=TEXT_BODY, line_spacing_pt=20)

    # Bottom callout
    add_shape_fill(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6),
                   RGBColor(0xEC, 0xFD, 0xF5), corner_radius=0.06)
    add_text_box(slide, Inches(1.1), Inches(6.55), Inches(11.1), Inches(0.5),
                 "111,360 frames processed across 3 full Hardpoint matches \u2014 all from raw broadcast footage",
                 font_size=16, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


def slide_features(prs):
    """Slide 5: Feature Engineering — concise two-column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "FEATURES", "12 Engineered Features from Two Data Sources")

    # Left: Scoreboard (6)
    add_shape_fill(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.6),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.4),
                 "SCOREBOARD (6 features)", font_size=18, color=ACCENT_BLUE, bold=True)

    score_feats = [
        "score_a, score_b \u2014 raw team scores",
        "score_diff \u2014 current lead/deficit",
        "score_rate_a, score_rate_b \u2014 scoring pace",
        "time_remaining_fraction \u2014 match progress",
    ]
    add_multiline_text(slide, Inches(1.1), Inches(3.0), Inches(5.0), Inches(2.5),
                       score_feats, font_size=16, color=TEXT_BODY, line_spacing_pt=26)

    # Right: Telemetry (6)
    add_shape_fill(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(3.6),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(7.1), Inches(2.4), Inches(5), Inches(0.4),
                 "MINIMAP TELEMETRY (6 features)", font_size=18, color=ACCENT_AMBER, bold=True)

    tele_feats = [
        "player_count_map, enemy_count_map \u2014 visible players",
        "player_advantage \u2014 friendly minus enemy",
        "30-second rolling averages for all three",
        "\u2014 smooths respawn cycles + fight dynamics",
    ]
    add_multiline_text(slide, Inches(7.1), Inches(3.0), Inches(5.0), Inches(2.5),
                       tele_feats, font_size=16, color=TEXT_BODY, line_spacing_pt=26)

    # Bottom callout
    add_shape_fill(slide, Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.8),
                   RGBColor(0xFF, 0xF7, 0xED), corner_radius=0.06)
    add_text_box(slide, Inches(2.7), Inches(6.3), Inches(7.9), Inches(0.6),
                 "Minimap telemetry features account for ~43% of model predictive power",
                 font_size=18, color=ACCENT_AMBER, bold=True, alignment=PP_ALIGN.CENTER)


def slide_model(prs):
    """Slide 6: Model — streamlined."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "MODEL", "XGBoost Classifier with Mirror Augmentation")

    # Left: Key concepts
    add_shape_fill(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.8),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.4),
                 "KEY DESIGN DECISIONS", font_size=18, color=ACCENT_BLUE, bold=True)

    add_multiline_text(slide, Inches(1.1), Inches(3.1), Inches(5.0), Inches(3.5), [
        "\u2022  XGBoost: 300 trees, depth 4, lr 0.05",
        "",
        "\u2022  Mirror augmentation: swap Team A \u2194 B",
        "    to double training data and balance classes",
        "",
        "\u2022  Hold-out validation: train on Matches 1 & 3,",
        "    test on Match 2 (different outcome, same video)",
    ], font_size=16, color=TEXT_BODY, line_spacing_pt=22)

    # Right: Hold-out design
    add_shape_fill(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.8),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(7.1), Inches(2.4), Inches(5), Inches(0.4),
                 "VALIDATION DESIGN", font_size=18, color=ACCENT_GREEN, bold=True)

    matches = [
        ("TRAIN", "Match 1: Team B wins (249\u2013151)", "38,100 frames", ACCENT_GREEN),
        ("TEST",  "Match 2: Team A wins (249\u2013166)", "34,320 frames", ACCENT_RED),
        ("TRAIN", "Match 3: Team A wins (250\u2013182)", "38,940 frames", ACCENT_GREEN),
    ]

    y_pos = Inches(3.2)
    for label, desc, frames, color in matches:
        # Badge
        add_shape_fill(slide, Inches(7.3), y_pos, Inches(1.0), Inches(0.55),
                       color, corner_radius=0.1)
        add_text_box(slide, Inches(7.3), y_pos + Inches(0.08),
                     Inches(1.0), Inches(0.4),
                     label, font_size=14, color=BG_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, Inches(8.5), y_pos + Inches(0.02),
                     Inches(3.5), Inches(0.3),
                     desc, font_size=16, color=TEXT_BLACK, bold=True)
        add_text_box(slide, Inches(8.5), y_pos + Inches(0.32),
                     Inches(3.5), Inches(0.25),
                     frames, font_size=13, color=TEXT_BODY)

        y_pos += Inches(0.9)

    add_text_box(slide, Inches(7.1), Inches(6.0), Inches(5.0), Inches(0.6),
                 "Different matches, different winners \u2014\na true generalization test.",
                 font_size=15, color=TEXT_BODY, line_spacing=21)


def slide_feature_importance(prs):
    """Slide 7: Feature Importance — top features only."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "RESULTS", "What Drives the Predictions?")

    # Top features only (keep it readable)
    features = [
        ("score_diff",           0.26, ACCENT_BLUE,  "Scoreboard"),
        ("player_alive_rolling", 0.16, ACCENT_AMBER, "Telemetry"),
        ("enemy_alive_rolling",  0.15, ACCENT_AMBER, "Telemetry"),
        ("advantage_rolling",    0.13, ACCENT_AMBER, "Telemetry"),
        ("score_rate_a",         0.07, ACCENT_BLUE,  "Scoreboard"),
        ("score_b",              0.05, ACCENT_BLUE,  "Scoreboard"),
    ]

    # Chart area
    add_shape_fill(slide, Inches(0.8), Inches(2.0), Inches(8.0), Inches(4.8),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(0.4),
                 "TOP 6 FEATURES BY IMPORTANCE", font_size=16, color=TEXT_BLACK, bold=True)

    max_bar_w = Inches(3.8)
    bar_h = Pt(22)
    start_y = Inches(2.9)
    row_h = Inches(0.55)

    for i, (name, imp, color, source) in enumerate(features):
        y = start_y + i * row_h

        # Feature name
        add_text_box(slide, Inches(1.1), y, Inches(2.8), Inches(0.4),
                     name, font_size=15, color=TEXT_BLACK, bold=True)

        # Bar
        bar_w = int(max_bar_w * (imp / 0.26))
        add_shape_fill(slide, Inches(4.2), y + Pt(5), bar_w, bar_h, color,
                       corner_radius=0.2)

        # Value on bar
        add_text_box(slide, Inches(4.2) + bar_w + Inches(0.12), y,
                     Inches(0.8), Inches(0.4),
                     f"{imp:.2f}", font_size=14, color=TEXT_BLACK, bold=True)

    # Right: summary cards
    add_stat_card(slide, Inches(9.2), Inches(2.2), Inches(3.3), Inches(1.4),
                  "~43%", "from minimap\ntelemetry", ACCENT_AMBER)

    add_stat_card(slide, Inches(9.2), Inches(3.9), Inches(3.3), Inches(1.4),
                  "~57%", "from scoreboard\nfeatures", ACCENT_BLUE)

    add_stat_card(slide, Inches(9.2), Inches(5.6), Inches(3.3), Inches(1.4),
                  "12", "total engineered\nfeatures", ACCENT_GREEN)


def slide_validation_chart(prs):
    """Slide 8: Validation chart — full width."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "VALIDATION", "Hold-Out Match: Predicted vs Actual")

    if os.path.exists(VALIDATION_CHART):
        slide.shapes.add_picture(
            VALIDATION_CHART, Inches(0.6), Inches(1.8),
            width=Inches(12.0), height=Inches(5.4),
        )


def slide_metrics(prs):
    """Slide 9: Key metrics — without 100% card."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "SUMMARY", "Pipeline Performance")

    # Top row — 3 cards
    cards_top = [
        ("111,360", "Total Frames\nProcessed", ACCENT_GREEN),
        ("99.3%", "Scoreboard Field\nCompleteness", ACCENT_BLUE),
        ("3 Matches", "Extracted from\nBroadcast Video", ACCENT_AMBER),
    ]

    card_w = Inches(3.6)
    gap = Inches(0.35)
    start_x = Inches(0.8)

    for i, (num, label, accent) in enumerate(cards_top):
        x = start_x + i * (card_w + gap)
        add_stat_card(slide, x, Inches(2.2), card_w, Inches(1.6), num, label, accent)

    # Bottom row — 3 cards
    cards_bot = [
        ("154,080", "Augmented\nTraining Rows", ACCENT_GREEN),
        ("43%", "Predictive Power\nfrom Telemetry", ACCENT_AMBER),
        ("59.94 fps", "Full Frame-Rate\nVideo Overlay", ACCENT_BLUE),
    ]

    for i, (num, label, accent) in enumerate(cards_bot):
        x = start_x + i * (card_w + gap)
        add_stat_card(slide, x, Inches(4.2), card_w, Inches(1.6), num, label, accent)

    # Bottom note
    add_shape_fill(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7),
                   BG_LIGHT, corner_radius=0.04)
    add_text_box(slide, Inches(1.1), Inches(6.3), Inches(11.1), Inches(0.5),
                 "The model correctly identifies the eventual winner for the majority of the held-out match,\n"
                 "demonstrating generalization across different matches and teams.",
                 font_size=15, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)


def slide_demo(prs):
    """Slide 10: Demo — with embedded video or thumbnail."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "DEMO", "Live Win Probability Overlay")

    # Try to embed the demo video thumbnail
    if os.path.exists(DEMO_THUMBNAIL):
        slide.shapes.add_picture(
            DEMO_THUMBNAIL, Inches(1.5), Inches(2.0),
            width=Inches(10.3), height=Inches(5.0),
        )
    else:
        # Placeholder box
        add_shape_fill(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(5.0),
                       BG_LIGHT, corner_radius=0.04)
        add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5),
                     "\u25B6  Play demo_clip.mp4",
                     font_size=36, color=TEXT_BLACK, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.5),
                     "30-second clip showing live win probability shifting during gameplay",
                     font_size=16, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)


def slide_contribution(prs):
    """Slide 11: My Contribution — emphasize student work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "DEVELOPMENT", "My Contribution")

    # Left: What I built
    add_shape_fill(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.8),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.4),
                 "WHAT I DESIGNED AND BUILT", font_size=18, color=ACCENT_GREEN, bold=True)

    add_multiline_text(slide, Inches(1.1), Inches(3.0), Inches(5.0), Inches(3.5), [
        "\u2022  Full pipeline architecture and design decisions",
        "\u2022  Template matching OCR approach (pivoted from EasyOCR)",
        "\u2022  YOLOv8 training and minimap detection pipeline",
        "\u2022  Feature engineering: 12 features from 2 data sources",
        "\u2022  Mirror augmentation strategy for class balancing",
        "\u2022  Hold-out validation methodology",
        "\u2022  Broadcast overlay design and rendering",
    ], font_size=16, color=TEXT_BODY, line_spacing_pt=24)

    # Right: Process & iteration
    add_shape_fill(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.8),
                   BG_LIGHT, corner_radius=0.04)

    add_text_box(slide, Inches(7.1), Inches(2.4), Inches(5), Inches(0.4),
                 "ITERATIVE DEVELOPMENT PROCESS", font_size=18, color=ACCENT_BLUE, bold=True)

    add_multiline_text(slide, Inches(7.1), Inches(3.0), Inches(5.0), Inches(3.5), [
        "\u2022  Started with EasyOCR, pivoted to template matching",
        "    after accuracy analysis on scoreboard fonts",
        "",
        "\u2022  Discovered minimap telemetry as a differentiator",
        "    (+43% predictive power vs. scoreboard alone)",
        "",
        "\u2022  Debugged frame offset recovery for video sync",
        "    across 454K frames of broadcast footage",
        "",
        "\u2022  All code, experiments, and analysis available",
    ], font_size=15, color=TEXT_BODY, line_spacing_pt=21)


def slide_references(prs):
    """Slide 12: References."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    slide_header(slide, "REFERENCES", "")

    refs = [
        "Anthropic. (2026). Claude Code [Large language model]. https://www.anthropic.com",
        "Bradski, G. (2000). The OpenCV Library. Dr. Dobb\u2019s Journal of Software Tools.",
        "Chen, T., & Guestrin, C. (2016). XGBoost: A Scalable Tree Boosting System. Proceedings of the 22nd ACM SIGKDD International Conference on Knowledge Discovery and Data Mining, 785\u2013794.",
        "Google. (2026). Gemini [Large language model]. https://gemini.google.com",
        "Jocher, G., Chaurasia, A., & Qiu, J. (2023). YOLO by Ultralytics (Version 8.0.0) [Computer software]. https://github.com/ultralytics/ultralytics",
        "Paszke, A., et al. (2019). PyTorch: An Imperative Style, High-Performance Deep Learning Library. Advances in Neural Information Processing Systems 32, 8024\u20138035.",
        "Settles, B. (2009). Active Learning Literature Survey. Computer Sciences Technical Report 1648, University of Wisconsin\u2013Madison.",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, ref in enumerate(refs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ref
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_BODY
        p.font.name = "Calibri"
        p.space_after = Pt(10)


def slide_closing(prs):
    """Slide 13: Thank you / Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_GREEN)

    add_text_box(slide, Inches(1.2), Inches(2.3), Inches(11), Inches(1.0),
                 "THANK YOU",
                 font_size=52, color=TEXT_BLACK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(3.6), Inches(2.3), ACCENT_GREEN, thickness=4)

    add_text_box(slide, Inches(1.2), Inches(4.1), Inches(11), Inches(0.6),
                 "Questions?",
                 font_size=28, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.5),
                 "Trevor Gregory",
                 font_size=20, color=TEXT_BLACK, bold=True,
                 alignment=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Generating slides...")

    slide_title(prs)                # 1. Title
    print("  [1/13] Title")

    slide_problem(prs)              # 2. Problem
    print("  [2/13] Problem")

    slide_pipeline(prs)             # 3. Pipeline
    print("  [3/13] Pipeline")

    slide_cv_extraction(prs)        # 4. CV Extraction
    print("  [4/13] Computer Vision")

    slide_features(prs)             # 5. Features
    print("  [5/13] Features")

    slide_model(prs)                # 6. Model
    print("  [6/13] Model")

    slide_feature_importance(prs)   # 7. Feature Importance
    print("  [7/13] Feature Importance")

    slide_validation_chart(prs)     # 8. Validation Chart
    print("  [8/13] Validation")

    slide_metrics(prs)              # 9. Metrics
    print("  [9/13] Metrics")

    slide_demo(prs)                 # 10. Demo
    print("  [10/13] Demo")

    slide_contribution(prs)         # 11. Contribution
    print("  [11/13] Contribution")

    slide_references(prs)           # 12. References
    print("  [12/13] References")

    slide_closing(prs)              # 13. Q&A
    print("  [13/13] Closing")

    prs.save(OUTPUT)
    size_mb = os.path.getsize(OUTPUT) / (1024 * 1024)
    print(f"\nDone! Saved \u2192 {OUTPUT} ({size_mb:.1f} MB)")
    print(f"Open with: open '{OUTPUT}'")


if __name__ == "__main__":
    main()
